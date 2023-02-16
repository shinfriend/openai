import os
import openai
from flask import Flask, redirect, render_template, request, url_for,Response
from pptx import Presentation
from pptx.util import Inches, Pt
import urllib.request
import time
import re
from docx import Document
from docx.shared import Inches
from pptx.enum.text import MSO_AUTO_SIZE   # MOS_AUTO_SIZEクラスのインポート
from pptx.enum.text import PP_ALIGN       # 段落の水平位置のEnume
from pptx.enum.dml import MSO_THEME_COLOR # テーマカラーのEnume
from pptx.dml.color import RGBColor
#from ppt2pdf import ppt2pdf

app = Flask(__name__)
openai.api_key = os.getenv("OPENAI_API_KEY")

@app.route("/", methods=("GET", "POST"))
def index():
    industry =""
    market = ""
    pestel_prompt = ""
    pestel_response = ""
    final_array=[]
    if request.method == "POST":       
        industry = request.form["industry"]
        market = request.form["market"]
        pestel = request.form["pestel"]
        swot = request.form["swot"]

        #Generate [PESEL Only] Report 
        if (pestel=='1' and swot=='0'):
            #Generate Prompt
            pestel_prompt = generate_prompt_pestel(industry,market,pestel)

            #Get Content from OpenAI 
            pestel_response = get_content(pestel_prompt)

            #Generate Power PowerPoint
            ppt_path = generate_powerpoint_pestel(pestel_response.choices[0].text,industry,market)
            doc_path = generate_word(pestel_response.choices[0].text,pestel,swot)
         

        elif (pestel=='0' and swot=='2'):
            #Generate Prompt
            swot_prompt = generate_prompt_swot(industry,market,swot)
            
            #Get Content from OpenAI 
            swot_response = get_content(swot_prompt)

            #Generate Power PowerPoint
            ppt_path = generate_powerpoint_swot(swot_response.choices[0].text,industry,market)
            doc_path = generate_word(swot_response.choices[0].text,pestel,swot)

        elif (pestel=='1' and swot=='2'):
            pestel_prompt = generate_prompt_pestel(industry,market,pestel)
            pestel_response = get_content(pestel_prompt)
            pestel_result = pestel_response.choices[0].text
            print(pestel_result)
            
            swot_prompt = generate_prompt_swot(industry,market,swot)
            swot_response = get_content(swot_prompt)
            swot_result = swot_response.choices[0].text
            print(swot_result)

            conclusion_prompt = generate_prompt_conclusion(industry,market)
            conclusion_response = get_content(conclusion_prompt)
            conclusion_result = conclusion_response.choices[0].text
            
            #Generate Power PowerPoint
            ppt_path = generate_powerpoint_pestel_swot(pestel_result,swot_result,conclusion_result,industry,market)
            
            doctxt = pestel_result + swot_result + conclusion_result
            doc_path = generate_word(doctxt,pestel,swot)


        else:
            print("This is default")       

        #Generate SWOT Prompt 

        #Generate PESTEL and SWOT Prompt 

        #Generate PESTE and SWOT Conclusion 

        #Set Report Path 
        final_array.append(ppt_path)
        final_array.append(doc_path)
        return redirect(url_for("index", result=final_array))  
    final_array1 = request.args.getlist("result")
   
    return render_template("index.html",final_array=final_array1)


def get_content(prompt_template):
    response = ""
    response= openai.Completion.create(
        model="text-davinci-003",
        prompt=prompt_template,  
        temperature=0.3,
        max_tokens=1000,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0
        )
    
    return response

def generate_powerpoint_pestel(result_text,industry,market):
    prs = Presentation()
    prs = Presentation('static/pestel_template.pptx')

    #Slide0 Conetents "Cover Slide" 
    generate_cover_contents(prs,industry,market,0)

    #Slide1 Conetents "Political" 
    generate_slide_contents(prs,result_text,"Political","Economic",1)

    #Slide2 Conetents "Economic" 
    generate_slide_contents(prs,result_text,"Economic","Social",2)

    #Slide3 Conetents "Social" 
    generate_slide_contents(prs,result_text,"Social","Technological",3)

    #Slide4 Conetents "Technological" 
    generate_slide_contents(prs,result_text,"Technological","Environmental",4)

   #Slide4 Conetents "Environmental" 
    generate_slide_contents(prs,result_text,"Environmental","Legal",5)

   #Slide5 Conetents "Legal" 
    generate_slide_contents(prs,result_text,"Legal"," ",6)

    demo_ppt = "static/presentation/pestel" +  time.strftime("%Y%m%d-%H%M%S") + ".pptx"
    prs.save(demo_ppt)
    #ppt2pdf(demo_ppt, 'static/presentation/pestel.pdf')
    return demo_ppt

def generate_powerpoint_swot(result_text,industry,market):
    prs = Presentation()
    prs = Presentation('static/swot_template.pptx')

    #Slide0 Conetents "Cover Slide" 
    generate_cover_contents(prs,industry,market,0)

    #Slide1 Conetents "Strength" 
    generate_slide_contents(prs,result_text,"Strengths","Weaknesses",1)

    #Slide2 Conetents "Weaknesses" 
    generate_slide_contents(prs,result_text,"Weaknesses","Opportunities",2)

    #Slide3 Conetents "Opportunities" 
    generate_slide_contents(prs,result_text,"Opportunities","Threats",3)

    #Slide4 Conetents "Threats" 
    generate_slide_contents(prs,result_text,"Threats","",4)

    demo_ppt = "static/presentation/swot" +  time.strftime("%Y%m%d-%H%M%S") + ".pptx"
    prs.save(demo_ppt)
    return demo_ppt

def generate_powerpoint_pestel_swot(pestel_result,swot_result,conclusion_result, industry,market):
    prs = Presentation()
    prs = Presentation('static/pestel_swot_template.pptx')

    #Slide1 Conetents "Cover Slide" 
    generate_cover_contents(prs,industry,market,0)

    #Slide2 Conetents "Political" 
    generate_slide_contents(prs,pestel_result,"Political","Economic",1)

    #Slide3 Conetents "Economic" 
    generate_slide_contents(prs,pestel_result,"Economic","Social",2)

    #Slide4 Conetents "Social" 
    generate_slide_contents(prs,pestel_result,"Social","Technological",3)

    #Slide5 Conetents "Technological" 
    generate_slide_contents(prs,pestel_result,"Technological","Environmental",4)

   #Slide6 Conetents "Environmental" 
    generate_slide_contents(prs,pestel_result,"Environmental","Legal",5)


   #Slide7 Conetents "Legal" 
    generate_slide_contents(prs,pestel_result,"Legal"," ",6)

    #SWOT####### 
    #Slide8 Conetents "Strength" 
    generate_slide_contents(prs,swot_result,"Strengths","Weaknesses",7)

    #Slide2 Conetents "Weaknesses" 
    generate_slide_contents(prs,swot_result,"Weaknesses","Opportunities",8)

    #Slide3 Conetents "Opportunities" 
    generate_slide_contents(prs,swot_result,"Opportunities","Threats",9)

    #Slide4 Conetents "Threats" 
    generate_slide_contents(prs,swot_result,"Threats","",10)

    #Conclusion & Recommendation
    generate_slide_contents_conclusion(prs,conclusion_result,"","",11)

    demo_ppt = "static/presentation/pestel_swot" +  time.strftime("%Y%m%d-%H%M%S") + ".pptx"
    prs.save(demo_ppt)
    return demo_ppt

def generate_slide_contents(pptxObj,txtResponse,strStar,strEnd,slideNo):

    #Remove [":"]
    strStartSearch = strStar + ":"
    strEndSearch = strEnd + ":"

    if strStartSearch in txtResponse:
        txtResponse = txtResponse.replace(strStartSearch,strStar)

    if strEndSearch in txtResponse:
        txtResponse = txtResponse.replace(strEndSearch,strEnd)   
    
    if "STRENGTHS" in txtResponse:
        txtResponse = txtResponse.replace("STRENGTHS",'Strengths')   

    if "WEAKNESSES" in txtResponse:
        txtResponse = txtResponse.replace("WEAKNESSES",'Weaknesses')   

    if "OPPORTUNITIES" in txtResponse:
        txtResponse = txtResponse.replace("OPPORTUNITIES",'Opportunities')   

    if "THREATS" in txtResponse:
        txtResponse = txtResponse.replace("THREATS",'Threats')           

    content_str = txtResponse
    #Add Two New Lines at the back [":"]
    strStstrStarartSearch = strStar + "\n\n"
    strEnd ="\n\n" + strEnd

    #Remove First [""]
    content_str= content_str[content_str.find(strStar)+len(strStar):content_str.rfind(strEnd)]
    #content_str = content_str.replace('•','')
    #content_str = content_str.replace(':', '',1)
    #content_str = content_str.replace('-', '')
    content_str = content_str.splitlines(True)

    print(content_str)

    while("\n" in content_str):
        content_str.remove("\n")
  
    print(content_str)
    content_str= ' '.join(content_str)

    sld1 = pptxObj.slides[slideNo]
    shapes = sld1.shapes

    for shape in sld1.shapes:           # スライド中の要素を抽出、種類を表示
        #print(shape.name)              # >> タイトル１、テキストボックス３、矢印：右４、吹き出し：円形１２
        if not shape.has_text_frame:    # shapeオブジェクトにTextFrameが含まれているか確認
            continue
        
    #textFrame = shape.text_frame    # 各種ShapeオブジェクトからTextFrameを取得   
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = content_str
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(18)
    font.color.rgb = RGBColor(255, 255, 255) 
    #tf.text = content_str
    tf.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.word_wrap = True    
    #tf.paragraphs[0].font.size = Pt(9)  # font size
    #tf.paragraphs[0].font.bold = False  # font bold
    #tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # font bold   

def generate_slide_contents_conclusion(pptxObj,txtResponse,strStar,strEnd,slideNo):

    sld1 = pptxObj.slides[slideNo]
    shapes = sld1.shapes

    for shape in sld1.shapes:           # スライド中の要素を抽出、種類を表示
        #print(shape.name)              # >> タイトル１、テキストボックス３、矢印：右４、吹き出し：円形１２
        if not shape.has_text_frame:    # shapeオブジェクトにTextFrameが含まれているか確認
            continue
        
    #textFrame = shape.text_frame    # 各種ShapeオブジェクトからTextFrameを取得   
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = txtResponse
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(18)
    font.color.rgb = RGBColor(255, 255, 255) 
    #tf.text = content_str
    tf.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.word_wrap = True    


def generate_cover_contents(pptxObj,industry, market, slideNo):
    
    txtCover= generate_cover_text(industry,market)
    
    sld1 = pptxObj.slides[slideNo]
    shapes = sld1.shapes

    for shape in sld1.shapes:           # スライド中の要素を抽出、種類を表示
        print(shape.name)              # >> タイトル１、テキストボックス３、矢印：右４、吹き出し：円形１２
        if not shape.has_text_frame:    # shapeオブジェクトにTextFrameが含まれているか確認
            continue
    
    txtCover = txtCover.replace('•','')
    #textFrame = shape.text_frame    # 各種ShapeオブジェクトからTextFrameを取得   
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = txtCover
    tf.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.word_wrap = True    
    #tf.paragraphs[0].font.size = Pt(28)  # font size
    tf.paragraphs[0].font.bold = False  # font bold
    tf.paragraphs[0].font.color.rgb = RGBColor(66, 205, 10)  # font bold   


def generate_word(result_text,pestel,swot):
    document = Document()
    demo_word=""
    if pestel=='1' and swot=='0':
        headerTxt = "PESTEL Analysis"
    elif pestel=='0' and swot=='2':
        headerTxt = "SWOT Analysis"
    else: 
        headerTxt = "PESTEL & SWOT Analysis with Conclusion"
    
    document.add_heading(headerTxt, 0)
    document.add_picture('static/united-logo.png', width=Inches(1.25))

    p = document.add_paragraph(result_text)
    document.add_page_break()
    demo_word = "static/presentation/word" +  time.strftime("%Y%m%d-%H%M%S") + ".docx"

    document.save(demo_word)
    return demo_word


def download_image(url, file_path, file_name):
    full_path = ""
    full_path = file_path + file_name + '.jpg'
    urllib.request.urlretrieve(url, full_path)
    return full_path

def generate_prompt(industry,market,pestel,swot):
    #print("This PESTEL value", pestel)
    #print("This SWOT value", swot)
    if (pestel=='1' and swot=='0'):
        print("Only PESTEL") 
    elif (pestel=='0' and swot=='2'):
        print("Only SWOT")
    elif (pestel=='1' and swot=='2'):
        print("Both PESTEL & SWOT")
    else:
        print("This is default")

    prompt_template = "Your task is to make a detailed PESTEL analysis for each section of the analysis for the XXXX industry and the markets under consideration is YYYY."
    #print(prompt_template)
    prompt_industry= re.sub("XXXX",industry,prompt_template)
    #print(prompt_industry)
    prompt_industry_market= re.sub("YYYY",market,prompt_industry)
    #print(prompt_industry_market)
    return prompt_industry_market

def generate_prompt_pestel(industry,market,pestel):
    prompt_pestel = ""
    if (pestel=='1'):
        prompt_template = "Your task is to make a detailed PESTEL analysis for each section of the analysis for the XXXX industry and the markets under consideration is specifically YYYY."
        prompt_industry= re.sub("XXXX",industry,prompt_template)
        prompt_pestel= re.sub("YYYY",market,prompt_industry)
        return prompt_pestel
    else:
        return prompt_pestel

def generate_prompt_swot(industry,market,swot):
    prompt_swot = ""
    if (swot=='2'):
        prompt_template = "Your task is to make a detailed SWOT analysis for each section of the analysis for the XXXX industry and the markets under consideration is specifically YYYY."
        prompt_industry= re.sub("XXXX",industry,prompt_template)
        prompt_swot= re.sub("YYYY",market,prompt_industry)
        return prompt_swot
    else:
        return prompt_swot

def generate_prompt_conclusion(industry,market):
    prompt_pestel_swot_con = ""
    prompt_template = "Can you make a synthesizing conclusion and recommendation based on the PESTEL and SWOT analysis for the XXXX industry and the markets under consideration is specifically YYYY?"
    prompt_industry= re.sub("XXXX",industry,prompt_template)
    prompt_pestel_swot_con= re.sub("YYYY",market,prompt_industry)
    return prompt_pestel_swot_con

def generate_cover_text(industry,market):
    prompt_template = "Analysis for the XXXX industry and the markets under consideration is YYYY."
    #print(prompt_template)
    prompt_industry= re.sub("XXXX",industry,prompt_template)
    #print(prompt_industry)
    prompt_industry_market= re.sub("YYYY",market,prompt_industry)
    #print(prompt_industry_market)
    return prompt_industry_market


@app.route('/progress')
def progress():
	def generate():
		x = 0
		
		while x <= 100:
			yield "data:" + str(x) + "\n\n"
			x = x + 10
			time.sleep(0.5)

	return Response(generate(), mimetype= 'text/event-stream')

